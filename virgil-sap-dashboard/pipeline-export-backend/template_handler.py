#!/usr/bin/env python3
"""
PowerPoint Template Handler using python-pptx

This module handles PowerPoint templates without programmatically extracting styles.
Instead, it uses templates as-is and overlays content on top of existing slide layouts.
"""

import os
import json
import tempfile
from typing import Dict, List, Optional, Any, Tuple
from pptx import Presentation
from pptx.slide import Slide
from pptx.shapes.base import BaseShape
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import logging

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class TemplateHandler:
    """
    Handles PowerPoint templates using python-pptx.
    Uses templates as-is without style extraction.
    """
    
    def __init__(self, template_path: str):
        """
        Initialize with a template file path.
        
        Args:
            template_path: Path to the .pptx template file
        """
        self.template_path = template_path
        self.presentation = None
        self.slide_layouts = []
        self.template_info = {}
        
    def load_template(self) -> bool:
        """
        Load the PowerPoint template.
        
        Returns:
            bool: True if successful, False otherwise
        """
        try:
            logger.info(f"Loading template: {self.template_path}")
            self.presentation = Presentation(self.template_path)
            
            # Extract basic template information
            self.template_info = {
                'slide_count': len(self.presentation.slides),
                'slide_layouts': len(self.presentation.slide_layouts),
                'slide_masters': len(self.presentation.slide_masters),
                'file_size': os.path.getsize(self.template_path) if os.path.exists(self.template_path) else 0
            }
            
            # Analyze available slide layouts
            self._analyze_slide_layouts()
            
            logger.info(f"Template loaded successfully: {self.template_info}")
            return True
            
        except Exception as e:
            logger.error(f"Failed to load template: {e}")
            return False
    
    def _analyze_slide_layouts(self):
        """Analyze available slide layouts in the template."""
        try:
            self.slide_layouts = []
            
            for i, layout in enumerate(self.presentation.slide_layouts):
                layout_type = self._get_layout_type(layout)
                layout_info = {
                    'index': i,
                    'name': layout.name or f"Layout {i}",
                    'type': layout_type,
                    'placeholders': len(layout.placeholders),
                    'shapes': len(layout.shapes)
                }
                self.slide_layouts.append(layout_info)
                logger.info(f"Layout {i}: {layout.name or f'Layout {i}'} - Type: {layout_type} - Placeholders: {len(layout.placeholders)}")
                
                # Debug: Log layout types distribution
                if i < 10:  # Only log first 10 layouts to avoid spam
                    logger.info(f"DEBUG: Layout {i} type '{layout_type}' for '{layout.name}'")
                
            logger.info(f"Found {len(self.slide_layouts)} slide layouts")
            
        except Exception as e:
            logger.error(f"Failed to analyze slide layouts: {e}")
    
    def _get_layout_type(self, layout) -> str:
        """Determine the type of slide layout."""
        try:
            # Check for common layout types based on placeholder content
            placeholders = layout.placeholders
            
            # Check for title slide
            if any(hasattr(p, 'placeholder_format') and p.placeholder_format.type == 1 for p in placeholders):  # Title
                return "title"
            
            # Check for content slide
            if any(hasattr(p, 'placeholder_format') and p.placeholder_format.type == 2 for p in placeholders):  # Content
                return "content"
            
            # Check for section header
            if any(hasattr(p, 'placeholder_format') and p.placeholder_format.type == 3 for p in placeholders):  # Section Header
                return "section_header"
            
            # Check for two content
            content_placeholders = [p for p in placeholders if hasattr(p, 'placeholder_format') and p.placeholder_format.type == 2]
            if len(content_placeholders) >= 2:
                return "two_content"
            
            # Fallback: try to determine by layout name
            layout_name = layout.name.lower() if layout.name else ""
            if "title" in layout_name:
                return "title"
            elif "content" in layout_name or "body" in layout_name:
                return "content"
            elif "section" in layout_name or "header" in layout_name:
                return "section_header"
            elif "two" in layout_name or "split" in layout_name:
                return "two_content"
            
            return "custom"
            
        except Exception as e:
            # Only log the error if it's not the expected attribute error
            if "placeholder_type" not in str(e):
                logger.error(f"Failed to determine layout type: {e}")
            # Don't log the placeholder_type error as it's expected for some layouts
            
            # Fallback based on layout name
            try:
                layout_name = layout.name.lower() if layout.name else ""
                if "title" in layout_name:
                    return "title"
                elif "content" in layout_name:
                    return "content"
                else:
                    return "content"  # Default to content
            except:
                return "content"  # Final fallback
    
    def create_slide_from_layout(self, layout_index: int, content: Dict[str, Any]) -> Optional[Slide]:
        """
        Create a new slide using a specific layout from the template.
        
        Args:
            layout_index: Index of the layout to use
            content: Dictionary containing slide content
            
        Returns:
            Slide: New slide object or None if failed
        """
        try:
            if not self.presentation:
                logger.error("Template not loaded")
                return None
            
            if layout_index >= len(self.presentation.slide_layouts):
                logger.error(f"Layout index {layout_index} out of range")
                return None
            
            # Create new slide using the specified layout
            slide_layout = self.presentation.slide_layouts[layout_index]
            slide = self.presentation.slides.add_slide(slide_layout)
            
            # Fill in content based on placeholders
            self._fill_slide_content(slide, content)
            
            logger.info(f"Created slide using layout {layout_index}")
            return slide
            
        except Exception as e:
            logger.error(f"Failed to create slide: {e}")
            return None
    
    def _fill_slide_content(self, slide: Slide, content: Dict[str, Any]):
        """
        Fill slide content based on available placeholders.
        
        Args:
            slide: The slide to fill
            content: Dictionary containing content to fill
        """
        try:
            # Get content fields
            title = content.get('title', '')
            body_text = content.get('content', '')
            subtitle = content.get('subtitle', '')
            
            logger.info(f"Filling slide content - Title: '{title[:50]}...', Body: '{body_text[:50]}...'")
            
            # Track what we've filled
            title_filled = False
            content_filled = False
            
            # Fill placeholders with specific types first
            for shape in slide.shapes:
                if hasattr(shape, 'placeholder_format'):
                    placeholder_type = shape.placeholder_format.type
                    
                    if placeholder_type == 1:  # Title
                        if hasattr(shape, 'text_frame'):
                            shape.text_frame.text = title
                            title_filled = True
                            logger.info(f"Filled title placeholder with: '{title[:30]}...'")
                    elif placeholder_type == 2:  # Content
                        if hasattr(shape, 'text_frame'):
                            shape.text_frame.text = body_text
                            content_filled = True
                            logger.info(f"Filled content placeholder with: '{body_text[:30]}...'")
                    elif placeholder_type == 3:  # Section Header
                        if hasattr(shape, 'text_frame'):
                            shape.text_frame.text = subtitle or title
                            title_filled = True
                            logger.info(f"Filled section header placeholder with: '{(subtitle or title)[:30]}...'")
            
            # If we haven't filled title or content, try to fill any text shape
            if not title_filled or not content_filled:
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame'):
                        current_text = shape.text_frame.text.strip()
                        
                        # Fill empty text frames or frames with placeholder text
                        if not current_text or current_text.lower() in ['click to edit master title style', 'click to edit master subtitle style', 'click to add text']:
                            if not title_filled:
                                shape.text_frame.text = title
                                title_filled = True
                                logger.info(f"Filled empty text frame with title: '{title[:30]}...'")
                            elif not content_filled and body_text:
                                shape.text_frame.text = body_text
                                content_filled = True
                                logger.info(f"Filled empty text frame with content: '{body_text[:30]}...'")
            
            # Final fallback: fill any remaining text frames
            if not title_filled or not content_filled:
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame') and shape.text_frame.text.strip() == '':
                        if not title_filled:
                            shape.text_frame.text = title
                            title_filled = True
                            logger.info(f"FINAL FALLBACK: Filled with title: '{title[:30]}...'")
                        elif not content_filled and body_text:
                            shape.text_frame.text = body_text
                            content_filled = True
                            logger.info(f"FINAL FALLBACK: Filled with content: '{body_text[:30]}...'")
            
            logger.info(f"Content filling complete - Title filled: {title_filled}, Content filled: {content_filled}")
                                
        except Exception as e:
            logger.error(f"Failed to fill slide content: {e}")
    
    def get_template_info(self) -> Dict[str, Any]:
        """
        Get template information.
        
        Returns:
            Dict containing template metadata
        """
        return self.template_info
    
    def get_available_layouts(self) -> List[Dict[str, Any]]:
        """
        Get available slide layouts.
        
        Returns:
            List of layout information dictionaries
        """
        return self.slide_layouts
    
    def save_presentation(self, output_path: str) -> bool:
        """
        Save the presentation to a file.
        
        Args:
            output_path: Path where to save the presentation
            
        Returns:
            bool: True if successful, False otherwise
        """
        try:
            if not self.presentation:
                logger.error("No presentation to save")
                return False
            
            self.presentation.save(output_path)
            logger.info(f"Presentation saved to: {output_path}")
            return True
            
        except Exception as e:
            logger.error(f"Failed to save presentation: {e}")
            return False
    
    def create_presentation_from_slides(self, slides_data: List[Dict[str, Any]], 
                                      output_path: str) -> bool:
        """
        Create a complete presentation from slide data.
        
        Args:
            slides_data: List of slide data dictionaries
            output_path: Path where to save the presentation
            
        Returns:
            bool: True if successful, False otherwise
        """
        try:
            if not self.presentation:
                logger.error("Template not loaded")
                return False
            
            # Clear existing slides (keep the template structure)
            # Note: We'll create new slides using the template layouts
            new_presentation = Presentation(self.template_path)
            
            # Remove any existing slides (keep only the template structure)
            while len(new_presentation.slides) > 0:
                rId = new_presentation.slides._sldIdLst[0].rId
                new_presentation.part.drop_rel(rId)
                new_presentation.slides._sldIdLst.remove(new_presentation.slides._sldIdLst[0])
            
            # Create slides based on data
            for i, slide_data in enumerate(slides_data):
                slide_type = slide_data.get('type', 'content')
                # Determine appropriate layout based on slide type and position
                layout_index = self._get_appropriate_layout(slide_type, i)
                
                # Log layout selection for debugging
                layout_name = new_presentation.slide_layouts[layout_index].name if layout_index < len(new_presentation.slide_layouts) else "Unknown"
                logger.info(f"Slide {i+1} ({slide_type}) using layout {layout_index} ({layout_name})")
                
                # Debug: Log what layouts were available for this slide type
                available_layouts = [l for l in self.slide_layouts if l['type'] in ['content', 'two_content', 'section_header', 'title']]
                logger.info(f"DEBUG: Available layouts for slide type '{slide_type}': {[l['index'] for l in available_layouts[:5]]}...")
                
                # Create slide
                slide_layout = new_presentation.slide_layouts[layout_index]
                slide = new_presentation.slides.add_slide(slide_layout)
                
                # Fill content
                self._fill_slide_content(slide, slide_data)
            
            # Save the presentation
            new_presentation.save(output_path)
            logger.info(f"Created presentation with {len(slides_data)} slides")
            return True
            
        except Exception as e:
            logger.error(f"Failed to create presentation: {e}")
            return False
    
    def _get_appropriate_layout(self, slide_type: str, slide_index: int = 0) -> int:
        """
        Get the most appropriate layout index for a slide type.
        
        Args:
            slide_type: Type of slide (title, content, etc.)
            slide_index: Index of the slide for variety
            
        Returns:
            int: Layout index to use
        """
        # Map slide types to layout types with more variety
        type_mapping = {
            'title': 'title',
            'executive_summary': 'section_header',
            'current_state': 'content',
            'business_challenges': 'two_content',
            'recommended_solutions': 'two_content',
            'solution_details': 'content',
            'benefits_roi': 'content',
            'implementation_roadmap': 'content',
            'investment_summary': 'content',
            'next_steps': 'section_header'
        }
        
        target_layout_type = type_mapping.get(slide_type, 'content')
        
        # Find layouts that have content placeholders (type 2)
        content_layouts = []
        for i, layout_info in enumerate(self.slide_layouts):
            # Check if this layout has content placeholders
            layout = self.presentation.slide_layouts[i]
            has_content_placeholder = False
            for shape in layout.shapes:
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 2:
                    has_content_placeholder = True
                    break
            
            if has_content_placeholder:
                content_layouts.append(i)
        
        logger.info(f"Found {len(content_layouts)} layouts with content placeholders: {content_layouts}")
        
        # For title slides, use layout 0 (usually the title layout)
        if slide_type == 'title':
            selected_layout = 0
            logger.info(f"TITLE SLIDE: Using layout 0 ({self.slide_layouts[0]['name']}) for slide 1")
            return selected_layout
        
        # For non-title slides, prioritize layouts with content placeholders
        if content_layouts:
            # Use progressive selection from content layouts
            layout_index = slide_index % len(content_layouts)
            selected_layout = content_layouts[layout_index]
            logger.info(f"CONTENT LAYOUT: Slide {slide_index + 1} using layout {selected_layout} ({self.slide_layouts[selected_layout]['name']}) for type '{slide_type}'")
            return selected_layout
        
        # Fallback: use progressive distribution from all layouts
        all_layouts = list(range(len(self.slide_layouts)))
        if len(all_layouts) > 1:
            # Use slide_index + 1 to start from layout 1 (skip title layout)
            layout_index = (slide_index + 1) % len(all_layouts)
            selected_layout = all_layouts[layout_index]
            logger.info(f"PROGRESSIVE FALLBACK: Slide {slide_index + 1} using layout {selected_layout} ({self.slide_layouts[selected_layout]['name']}) for type '{slide_type}'")
            return selected_layout
        
        # Final fallback to first available layout
        logger.info(f"FINAL FALLBACK: Using layout 0 for slide type '{slide_type}'")
        return 0 if self.slide_layouts else 0


def process_template_request(template_path: str, slides_data: List[Dict[str, Any]], 
                           output_path: str) -> Dict[str, Any]:
    """
    Process a template request and create a presentation.
    
    Args:
        template_path: Path to the template file
        slides_data: List of slide data
        output_path: Path where to save the output
        
    Returns:
        Dict containing success status and metadata
    """
    try:
        # Initialize template handler
        handler = TemplateHandler(template_path)
        
        # Load template
        if not handler.load_template():
            return {
                'success': False,
                'error': 'Failed to load template'
            }
        
        # Get template info
        template_info = handler.get_template_info()
        available_layouts = handler.get_available_layouts()
        
        # Create presentation
        success = handler.create_presentation_from_slides(slides_data, output_path)
        
        if success:
            return {
                'success': True,
                'template_info': template_info,
                'available_layouts': available_layouts,
                'slides_created': len(slides_data),
                'output_path': output_path
            }
        else:
            return {
                'success': False,
                'error': 'Failed to create presentation'
            }
            
    except Exception as e:
        logger.error(f"Template processing failed: {e}")
        return {
            'success': False,
            'error': str(e)
        }


if __name__ == "__main__":
    # Example usage
    template_path = "example_template.pptx"
    slides_data = [
        {
            'type': 'title',
            'title': 'Sample Presentation',
            'content': 'This is a sample presentation created with python-pptx'
        },
        {
            'type': 'content',
            'title': 'Slide 1',
            'content': 'This is the content of slide 1'
        }
    ]
    
    result = process_template_request(template_path, slides_data, "output.pptx")
    print(json.dumps(result, indent=2)) 